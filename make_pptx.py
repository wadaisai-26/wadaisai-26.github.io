from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ---- color palette ----
C_NAVY   = RGBColor(0x0f, 0x34, 0x60)
C_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
C_RED    = RGBColor(0xe6, 0x39, 0x46)
C_GOLD   = RGBColor(0xff, 0xbd, 0x59)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_LGRAY  = RGBColor(0xfa, 0xfa, 0xfa)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_YELLOW = RGBColor(0xff, 0xf8, 0xe1)
C_YTEXT  = RGBColor(0x85, 0x64, 0x04)

# ---- helpers ----
def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_text_box(slide, lines, x, y, w, h,
                 bg=None, line_color=None,
                 padding=Inches(0.2)):
    rect = add_rect(slide, x, y, w, h, fill=bg, line=line_color)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left   = padding
    tf.margin_right  = padding
    tf.margin_top    = padding
    tf.margin_bottom = padding
    first = True
    for (text, size, bold, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return rect

# ============================================================
# SLIDE 1 — 表紙
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# background
add_rect(sl, 0, 0, W, H, fill=C_DARK)
# accent circle top-right
circle = sl.shapes.add_shape(9, W - Inches(4), -Inches(1.5), Inches(5), Inches(5))
circle.fill.solid(); circle.fill.fore_color.rgb = RGBColor(0xe6,0x39,0x46)
circle.line.fill.background()
# use a muted red for the accent circle
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x3a, 0x18, 0x22)

# red tag
add_text_box(sl,
    [("許可申請資料 2026", 11, True, C_WHITE, PP_ALIGN.CENTER)],
    Inches(4.66), Inches(1.4), Inches(1.5), Inches(0.35), bg=C_RED)

# title
add_text(sl, "こどもまつり2026", Inches(1.5), Inches(2.0), Inches(10), Inches(0.9),
         size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "キッチンカー招致のご承認依頼", Inches(1.5), Inches(2.85), Inches(10), Inches(0.9),
         size=36, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
add_text(sl, "安全・安心なフードサービスの導入について",
         Inches(1.5), Inches(3.75), Inches(10), Inches(0.5),
         size=16, color=RGBColor(0xcc,0xcc,0xcc), align=PP_ALIGN.CENTER)

# divider
add_rect(sl, Inches(2), Inches(4.4), Inches(9.33), Pt(1),
         fill=RGBColor(0x44,0x44,0x55))

# meta
add_text(sl, "📅 開催予定：2026年度　　🏫 主催：○○大学 こどもまつり実行委員会",
         Inches(1.5), Inches(4.55), Inches(10), Inches(0.5),
         size=13, color=RGBColor(0x99,0x99,0xbb), align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — なぜキッチンカーが必要か
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, 0, 0, W, H, fill=C_LGRAY)

# header
add_text(sl, "REASON", Inches(0), Inches(0.55), W, Inches(0.35),
         size=11, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
add_text(sl, "なぜキッチンカーが必要か", Inches(0), Inches(0.85), W, Inches(0.65),
         size=30, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(5.66), Inches(1.52), Inches(2), Pt(3), fill=C_RED)

# Card 1 - 熱中症
card1_x = Inches(0.6)
card_y   = Inches(2.0)
card_w   = Inches(5.8)
card_h   = Inches(4.5)
add_rect(sl, card1_x, card_y, card_w, card_h, fill=C_WHITE)
add_rect(sl, card1_x, card_y, card_w, Pt(5), fill=C_RED)

add_text(sl, "🌞", card1_x + Inches(0.3), card_y + Inches(0.25),
         Inches(1), Inches(0.6), size=32)
add_text(sl, "熱中症リスクへの対応",
         card1_x + Inches(0.3), card_y + Inches(0.9), card_w - Inches(0.5), Inches(0.5),
         size=18, bold=True, color=C_DARK)
add_text(sl,
    "こどもまつりは屋外・夏季開催のため、長時間の活動により子どもたちの熱中症リスクが高まります。\n\n"
    "現状、昼食を提供できる場所が大学内にないため、参加者が学外へ出る必要があり、"
    "炎天下での移動が危険を招きます。",
    card1_x + Inches(0.3), card_y + Inches(1.45), card_w - Inches(0.5), Inches(1.8),
    size=13, color=C_GRAY)
add_text_box(sl,
    [("⚠ 食事・水分補給の場を確保することが安全管理上の急務です", 12, True, C_YTEXT, PP_ALIGN.LEFT)],
    card1_x + Inches(0.3), card_y + Inches(3.3), card_w - Inches(0.55), Inches(0.55),
    bg=C_YELLOW, padding=Inches(0.12))

# Card 2 - 滞在促進
card2_x = Inches(6.9)
add_rect(sl, card2_x, card_y, card_w, card_h, fill=C_WHITE)
add_rect(sl, card2_x, card_y, card_w, Pt(5), fill=C_NAVY)

add_text(sl, "🎪", card2_x + Inches(0.3), card_y + Inches(0.25),
         Inches(1), Inches(0.6), size=32)
add_text(sl, "来場者の滞在促進",
         card2_x + Inches(0.3), card_y + Inches(0.9), card_w - Inches(0.5), Inches(0.5),
         size=18, bold=True, color=C_DARK)
add_text(sl,
    "昼食場所がないことで、来場者が昼をはさんで一時退場・早退してしまう問題があります。\n\n"
    "キッチンカーの設置により会場内で食事・休憩できる環境を作り、"
    "イベントを通じた体験の充実と来場者満足度の向上を図ります。",
    card2_x + Inches(0.3), card_y + Inches(1.45), card_w - Inches(0.5), Inches(1.8),
    size=13, color=C_GRAY)
add_text_box(sl,
    [("✓ 来場者の滞在時間延長 → イベント全体の活性化", 12, True, RGBColor(0x1a,0x73,0xe8), PP_ALIGN.LEFT)],
    card2_x + Inches(0.3), card_y + Inches(3.3), card_w - Inches(0.55), Inches(0.55),
    bg=RGBColor(0xe8,0xf0,0xfe), padding=Inches(0.12))

# ============================================================
# SLIDE 3 — モビマルの信頼性
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, 0, 0, W, H, fill=C_NAVY)

add_text(sl, "TRUST & SCALE", Inches(0), Inches(0.45), W, Inches(0.35),
         size=11, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
add_text(sl, "利用するサービス「モビマル」の信頼性・規模",
         Inches(0), Inches(0.78), W, Inches(0.6),
         size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# stats
stats = [
    ("16,000件+", "イベント派遣実績"),
    ("5,400台+",  "登録キッチンカー数"),
    ("No.1",      "掲載台数\n主要6社中"),
]
sx = Inches(0.5)
sw = Inches(3.9)
for i, (num, label) in enumerate(stats):
    x = sx + i * (sw + Inches(0.2))
    add_rect(sl, x, Inches(1.65), sw, Inches(1.6),
             fill=RGBColor(0x1e,0x3d,0x6e), line=RGBColor(0x44,0x66,0x99))
    add_text(sl, num, x, Inches(1.75), sw, Inches(0.8),
             size=34, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text(sl, label, x, Inches(2.5), sw, Inches(0.65),
             size=13, color=RGBColor(0xcc,0xcc,0xff), align=PP_ALIGN.CENTER)

# badges
badges = [
    ("📈 東証プライム上場企業が運営",
     "株式会社シンクロ・フードが事業運営。財務・コンプライアンス体制が整備された信頼ある運営母体。"),
    ("🏛 一般社団法人 日本移動販売協会",
     "2019年設立の業界団体が後援。業界標準の衛生・安全基準を推進。"),
    ("🔍 審査済みの事業者のみ登録",
     "許可証・保険加入を確認した事業者がマッチング。素性不明の業者が入ることがない仕組み。"),
]
bw = Inches(4.0)
by = Inches(3.55)
for i, (title, body) in enumerate(badges):
    x = Inches(0.4) + i * (bw + Inches(0.26))
    add_rect(sl, x, by, bw, Inches(2.8),
             fill=RGBColor(0x16,0x2d,0x55), line=RGBColor(0x33,0x55,0x88))
    add_text(sl, title, x + Inches(0.2), by + Inches(0.2), bw - Inches(0.35), Inches(0.5),
             size=13, bold=True, color=C_WHITE)
    add_text(sl, body, x + Inches(0.2), by + Inches(0.75), bw - Inches(0.35), Inches(1.8),
             size=12, color=RGBColor(0xbb,0xcc,0xee))

# ============================================================
# SLIDE 4 — 許可証の概要
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, 0, 0, W, H, fill=C_LGRAY)

add_text(sl, "PERMITS", Inches(0), Inches(0.45), W, Inches(0.35),
         size=11, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
add_text(sl, "キッチンカー事業者が取得する許可証",
         Inches(0), Inches(0.78), W, Inches(0.6),
         size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

permits = [
    ("📄", "食品営業許可証",
     "食品衛生法に基づき、各都道府県の保健所が発行。調理設備・衛生設備が基準を満たすと認定された"
     "車両のみ取得可能。当日は必ず原本を携帯・掲示します。",
     "発行：都道府県保健所"),
    ("🧑‍🍳", "食品衛生責任者",
     "各営業車に1名以上の有資格者配置が義務付けられています。食品衛生に関する知識と管理責任を"
     "持った人員が常駐し、当日の衛生管理を担います。",
     "修了証：食品衛生協会"),
    ("✅", "HACCP準拠の衛生管理",
     "2021年6月より全食品事業者に義務化されたHACCPに基づく衛生管理を実施。仕入れ・調理・"
     "提供の各工程で危害リスクを分析・管理します。",
     "根拠：改正食品衛生法"),
]
pw = Inches(4.0)
for i, (icon, title, body, tag) in enumerate(permits):
    x = Inches(0.4) + i * (pw + Inches(0.26))
    add_rect(sl, x, Inches(1.7), pw, Inches(3.3), fill=C_WHITE)
    add_rect(sl, x, Inches(1.7), pw, Pt(4), fill=C_RED)
    add_text(sl, icon, x + Inches(0.2), Inches(1.85), pw, Inches(0.6), size=28)
    add_text(sl, title, x + Inches(0.2), Inches(2.5), pw - Inches(0.3), Inches(0.45),
             size=15, bold=True, color=C_DARK)
    add_text(sl, body, x + Inches(0.2), Inches(2.95), pw - Inches(0.3), Inches(1.5),
             size=12, color=C_GRAY)
    add_text_box(sl,
        [(tag, 11, True, C_NAVY, PP_ALIGN.CENTER)],
        x + Inches(0.2), Inches(4.6), pw - Inches(0.35), Inches(0.3),
        bg=RGBColor(0xf0,0xf4,0xff), padding=Inches(0.06))

# liability note
add_rect(sl, Inches(0.4), Inches(5.2), W - Inches(0.8), Inches(1.7),
         fill=C_YELLOW, line=C_GOLD)
add_text(sl, "⚖ 法的責任の所在について",
         Inches(0.65), Inches(5.35), W - Inches(1.2), Inches(0.38),
         size=14, bold=True, color=C_DARK)
add_text(sl,
    "食品衛生法上、食中毒などの事故が発生した場合の営業上の責任は、許可を受けたキッチンカー事業者が負います。"
    "大学は食品を製造・販売する立場ではなく、場所を提供するにとどまるため、食中毒に関する直接的な法的責任を負う立場にはありません。",
    Inches(0.65), Inches(5.72), W - Inches(1.2), Inches(0.9),
    size=13, color=RGBColor(0x44,0x33,0x00))

# ============================================================
# SLIDE 5 — PL保険
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, 0, 0, W, H, fill=C_DARK)

add_text(sl, "PL INSURANCE", Inches(0), Inches(0.45), W, Inches(0.35),
         size=11, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
add_text(sl, "PL保険（製造物責任保険）の概要",
         Inches(0), Inches(0.78), W, Inches(0.6),
         size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# left panel
add_rect(sl, Inches(0.4), Inches(1.65), Inches(8.0), Inches(5.5),
         fill=RGBColor(0x22,0x22,0x44), line=RGBColor(0x44,0x44,0x66))
add_text(sl, "補償対象となる事故",
         Inches(0.7), Inches(1.85), Inches(7.4), Inches(0.45),
         size=16, bold=True, color=C_GOLD)

items = [
    ("✓ 食中毒",
     "販売した飲食物による食中毒発症時の治療費・入院費・休業損害・慰謝料"),
    ("✓ 異物混入による健康被害",
     "ドリンク・食品への異物混入による歯欠損などの身体的損害"),
    ("✓ 物損事故",
     "提供時の飲料こぼし等による衣類・所持品への損害"),
    ("✓ 施設賠償責任保険（セット）",
     "看板落下など設備由来の事故による来場者のケガ"),
]
for j, (t, b) in enumerate(items):
    y = Inches(2.45) + j * Inches(1.05)
    add_rect(sl, Inches(0.65), y, Inches(0.3), Inches(0.3),
             fill=C_RED)
    add_text(sl, t, Inches(1.05), y - Inches(0.03), Inches(7.0), Inches(0.38),
             size=13, bold=True, color=C_WHITE)
    add_text(sl, b, Inches(1.05), y + Inches(0.32), Inches(7.0), Inches(0.45),
             size=12, color=RGBColor(0xbb,0xcc,0xee))

# right panel — amount
add_rect(sl, Inches(8.7), Inches(1.65), Inches(4.2), Inches(2.4), fill=C_RED)
add_text(sl, "モビマル加入の補償上限額",
         Inches(8.7), Inches(1.85), Inches(4.2), Inches(0.4),
         size=12, color=RGBColor(0xff,0xdd,0xdd), align=PP_ALIGN.CENTER)
add_text(sl, "3億円",
         Inches(8.7), Inches(2.25), Inches(4.2), Inches(0.9),
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "（PL保険 ＋ 施設賠償責任保険）",
         Inches(8.7), Inches(3.12), Inches(4.2), Inches(0.35),
         size=12, color=RGBColor(0xff,0xdd,0xdd), align=PP_ALIGN.CENTER)

# right panel — univ note
add_rect(sl, Inches(8.7), Inches(4.2), Inches(4.2), Inches(2.95),
         fill=RGBColor(0x22,0x22,0x44), line=RGBColor(0x44,0x44,0x66))
add_text(sl, "🏫 大学への影響について",
         Inches(8.9), Inches(4.38), Inches(4.0), Inches(0.4),
         size=13, bold=True, color=C_GOLD)
add_text(sl,
    "PL保険はキッチンカー事業者が加入・負担するものです。"
    "万一の事故も事業者の保険で対応されるため、大学が賠償責任を肩代わりする状況は生じません。\n\n"
    "事前に保険証書のコピーを提出いただくことで書面での確認も可能です。",
    Inches(8.9), Inches(4.85), Inches(3.8), Inches(2.1),
    size=12, color=RGBColor(0xbb,0xcc,0xee))

# ============================================================
# SLIDE 6 — まとめ
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
# gradient-ish: navy left, red right
add_rect(sl, 0, 0, W, H, fill=C_NAVY)
add_rect(sl, W // 2, 0, W // 2, H, fill=RGBColor(0xc0, 0x28, 0x33))

add_text(sl, "まとめ・ご承認のお願い",
         Inches(0), Inches(0.3), W, Inches(0.65),
         size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "以上の点から、安全・安心なキッチンカー招致の実現をお願い申し上げます",
         Inches(0), Inches(0.95), W, Inches(0.4),
         size=14, color=RGBColor(0xdd,0xdd,0xff), align=PP_ALIGN.CENTER)

points = [
    ("1", "必要性：熱中症対策と来場者サービス",
     "昼食場所の確保は子どもたちの安全管理上の急務であり、滞在促進によるイベント充実にも貢献します"),
    ("2", "信頼性：東証プライム上場企業運営・実績16,000件超",
     "モビマルは業界最大規模のマッチングサービスであり、審査済み事業者のみが登録されています"),
    ("3", "安全管理：法的許可証・HACCP・最大3億円のPL保険",
     "食品営業許可証・食品衛生責任者・PL保険はすべてキッチンカー事業者が取得・加入しており、食中毒等の法的責任は事業者が負います"),
    ("4", "大学の責任範囲：場所の提供に限定",
     "大学は食品の製造・販売主体ではないため、食品衛生法上の直接的責任は生じません。事業者との覚書締結で責任関係を明文化します"),
]
for i, (num, title, body) in enumerate(points):
    y = Inches(1.55) + i * Inches(1.2)
    add_rect(sl, Inches(0.4), y + Inches(0.05),
             Inches(0.5), Inches(0.5), fill=C_WHITE)
    add_text(sl, num, Inches(0.4), y, Inches(0.5), Inches(0.6),
             size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, title, Inches(1.05), y, Inches(11.5), Inches(0.42),
             size=14, bold=True, color=C_WHITE)
    add_text(sl, body, Inches(1.05), y + Inches(0.42), Inches(11.5), Inches(0.6),
             size=12, color=RGBColor(0xcc,0xdd,0xff))

# request box
add_rect(sl, Inches(3.5), Inches(6.5), Inches(6.33), Inches(0.8), fill=C_WHITE)
add_text(sl, "キッチンカー招致のご承認をよろしくお願い申し上げます 🙏",
         Inches(3.5), Inches(6.55), Inches(6.33), Inches(0.65),
         size=15, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

# ============================================================
out = r"C:\Users\gamep\OneDrive\デスクトップ\kodomomaturi\キッチンカー招致申請.pptx"
prs.save(out)
print("Saved:", out)
